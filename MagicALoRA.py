import os
import json
import logging
import tkinter as tk
from tkinter import filedialog, messagebox, ttk
from datetime import datetime
import threading
import fitz  # PyMuPDF
from pptx import Presentation
from moviepy.editor import VideoFileClip
import speech_recognition as sr
from pydub import AudioSegment
import csv
import ebooklib
from ebooklib import epub
from bs4 import BeautifulSoup
import xml.etree.ElementTree as ET
import zipfile
import pandas as pd
from pydub.silence import split_on_silence
from pytube import YouTube
import vlc
from docx import Document
import re
from difflib import SequenceMatcher
from tkinterdnd2 import TkinterDnD

# Temporary directory
temp_dir = 'temp'
if not os.path.exists(temp_dir):
    os.makedirs(temp_dir)

# Logger configuration
def configure_logger(log_name):
    timestamp = datetime.now().strftime("%Y%m%d-%H%M%S")
    log_file = os.path.join(temp_dir, f'{log_name}-{timestamp}.log')
    logging.basicConfig(filename=log_file, level=logging.INFO, format='%(asctime)s - %(levelname)s - %(message)s')
    return log_file

# Global variable for language
lang = {}

# Load language translations from XML file
def load_language(language_code):
    global lang
    try:
        tree = ET.parse(f'{language_code}.xml')
        root = tree.getroot()
        lang = {}
        for field in root.find('fields'):
            lang[field.tag] = field.text
        update_language(lang)
    except Exception as e:
        logging.error(f"Failed to load language file: {language_code}.xml - {str(e)}")
        messagebox.showerror("Error", f"Failed to load language file: {str(e)}")

# Function to update the interface language
def update_language(lang):
    try:
        tab_control.tab(0, text=lang['convert'])
        tab_control.tab(1, text=lang['testSRT'])
        tab_control.tab(2, text=lang['createTextFile'])
        tab_control.tab(3, text=lang['createJson'])
        tab_control.tab(4, text=lang['setup'])
        url_label.config(text=lang['orEnterUrl'])
        browse_button.config(text=lang['browseVideo'])
        start_button.config(text=lang['startProcessing'])
        download_audio_only_checkbox.config(text=lang['downloadAudioOnly'])
        url_label_video.config(text=lang['orEnterUrlVideo'])
        video_browse_button.config(text=lang['browseVideo'])
        srt_browse_button.config(text=lang['browseSRTFile'])
        load_video_button.config(text=lang['loadAndPlayVideo'])
        create_text_file_button.config(text=lang['startCreateTextFile'])
        create_json_button.config(text=lang['startCreateJson'])
        chapter_keywords_label.config(text=lang['enterChapterKeywords'])
        add_keyword_button.config(text=lang['addKeyword'])
        remove_keyword_button.config(text=lang['removeKeyword'])
        add_directory_button.config(text=lang['addDirectory'])
        remove_directory_button.config(text=lang['removeDirectory'])
        add_ignore_directory_button.config(text=lang['addIgnoreDirectory'])
        remove_ignore_directory_button.config(text=lang['removeIgnoreDirectory'])
        set_output_path_button.config(text=lang['setTextOutput'])
        set_json_output_path_button.config(text=lang['setJsonOutput'])
        process_subfolders_checkbox.config(text=lang['processSubfolders'])
        save_config_button.config(text=lang['saveConfiguration'])
        load_config_button.config(text=lang['loadConfiguration'])
        limit_search_label.config(text=lang['limitSearch'])
        limit_search_menu['values'] = [lang['noLimit'], lang['lastProducedPerType'], lang['lastProducedInFolder'], lang['lastProducedSimilarTitle']]
    except KeyError as e:
        logging.error(f"Missing language key: {str(e)}")
        messagebox.showerror("Error", f"Missing language key: {str(e)}")

# Functions to handle different file types
def remove_headers_footers(text):
    lines = text.split('\n')
    if len(lines) > 3:
        return '\n'.join(lines[1:-1])
    return text

def handle_text_file(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8', errors='replace') as file:
            text = file.read()
        return remove_headers_footers(text), file_path
    except Exception as e:
        logging.error(f"Failed to read or process text file: {file_path} - {str(e)}")
        return f"Failed to read or process text file: {str(e)}", None

def handle_pdf_file(file_path):
    try:
        doc = fitz.open(file_path)
        text = [page.get_text("text") for page in doc]
        doc.close()
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        logging.error(f"Failed to process PDF file: {file_path} - {str(e)}")
        return f"Failed to process PDF file: {file_path} - {str(e)}", None

def handle_word_file(file_path):
    try:
        doc = Document(file_path)
        text = '\n'.join([para.text for para in doc.paragraphs])
        return remove_headers_footers(text), file_path
    except Exception as e:
        logging.error(f"Failed to process Word file: {file_path} - {str(e)}")
        return f"Failed to process Word file: {file_path} - {str(e)}", None

def handle_ppt_file(file_path):
    try:
        ppt = Presentation(file_path)
        text = [shape.text for slide in ppt.slides for shape in slide.shapes if hasattr(shape, "text")]
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        logging.error(f"Failed to process PowerPoint file: {file_path} - {str(e)}")
        return f"Failed to process PowerPoint file: {file_path} - {str(e)}", None

def handle_excel_file(file_path):
    try:
        df = pd.read_excel(file_path)
        return df.to_csv(index=False), file_path
    except Exception as e:
        logging.error(f"Failed to process Excel file: {file_path} - {str(e)}")
        return f"Failed to process Excel file: {file_path} - {str(e)}", None

def handle_csv_file(file_path):
    try:
        with open(file_path, mode='r', encoding='utf-8', newline='') as f:
            reader = csv.reader(f)
            data = list(reader)
        return '\n'.join([','.join(row) for row in data]), file_path
    except Exception as e:
        logging.error(f"Failed to process CSV file: {file_path} - {str(e)}")
        return f"Failed to process CSV file: {file_path} - {str(e)}", None

def handle_epub_file(file_path):
    try:
        book = epub.read_epub(file_path)
        text = []
        for item in book.get_items():
            if item.get_type() == ebooklib.ITEM_DOCUMENT:
                soup = BeautifulSoup(item.get_content(), 'html.parser')
                text.append(soup.get_text())
        return remove_headers_footers('\n'.join(text)), file_path
    except Exception as e:
        logging.error(f"Failed to process EPUB file: {file_path} - {str(e)}")
        return f"Failed to process EPUB file: {file_path} - {str(e)}", None

def handle_xml_gan_file(file_path):
    try:
        tree = ET.parse(file_path)
        root = tree.getroot()
        texts = [elem.text for elem in root.iter() if elem.text is not None]
        return '\n'.join(texts), file_path
    except Exception as e:
        logging.error(f"Failed to process XML/GAN file: {file_path} - {str(e)}")
        return f"Failed to process XML/GAN file: {file_path} - {str(e)}", None

def handle_audio_file(file_path):
    if file_path.lower().endswith('.m4a'):
        sound = AudioSegment.from_file(file_path, format='m4a')
        wav_path = file_path.replace('.m4a', '.wav')
        sound.export(wav_path, format='wav')
        file_path = wav_path

    recognizer = sr.Recognizer()
    with sr.AudioFile(file_path) as source:
        audio_data = recognizer.record(source)
        try:
            return recognizer.recognize_google(audio_data, language='it-IT'), file_path
        except sr.UnknownValueError:
            logging.warning(f"Speech not understood in file: {file_path}")
            return "Speech not understood", file_path
        except sr.RequestError as e:
            logging.error(f"Speech recognition request failed for file: {file_path} - {e}")
            return f"Speech recognition request failed; {e}", file_path

def handle_generic_video_file(file_path):
    try:
        audio_path = extract_audio_from_video(file_path)
        text = transcribe_audio(audio_path)
        os.remove(audio_path)
        return text, file_path
    except Exception as e:
        logging.error(f"Failed to process video file: {file_path} - {str(e)}")
        return f"Failed to process video file: {file_path} - {str(e)}", None

def extract_audio_from_video(video_path):
    video = VideoFileClip(video_path)
    audio_path = os.path.join(temp_dir, "temp_audio.wav")
    video.audio.write_audiofile(audio_path)
    return audio_path

def transcribe_audio(audio_path, language='it-IT'):
    recognizer = sr.Recognizer()
    with sr.AudioFile(audio_path) as source:
        audio_data = recognizer.record(source)
        try:
            return recognizer.recognize_google(audio_data, language=language)
        except sr.UnknownValueError:
            return "Speech not understood"
        except sr.RequestError as e:
            return f"Could not request results; {e}"

def write_to_output(content, output_dir, file_index, original_path):
    output_file_path = os.path.join(output_dir, f'model_{file_index}.txt')
    with open(output_file_path, 'a', encoding='utf-8') as file:
        file.write(f"\nOriginal file path: {original_path}\nFile content:\n{content}\n")
    return file_index + 1

def handle_zip_file(zip_path):
    try:
        with zipfile.ZipFile(zip_path, 'r') as z:
            z.extractall(temp_dir)
            extracted_files = z.namelist()
            for file_name in extracted_files:
                internal_path = os.path.join(temp_dir, file_name)
                if os.path.isfile(internal_path):
                    content, _ = handle_file(internal_path)
                    if content and not content.startswith("Unsupported"):
                        return f"{content} (from {file_name} in {zip_path})", zip_path
                    os.remove(internal_path)
        return "No supported files found or failed to process", None
    except Exception as e:
        logging.error(f"Failed to process ZIP file: {zip_path} - {str(e)}")
        return f"Failed to process ZIP file: {str(e)}", None

def explore_directory(directory, output_dir, ignore_dirs, process_subfolders, limit_search):
    file_index = 1
    for root, dirs, files in os.walk(directory):
        dirs[:] = [d for d in dirs if os.path.join(root, d) not in ignore_dirs]
        if not process_subfolders:
            dirs[:] = []
        files = limit_files_search(files, limit_search)
        for file in files:
            file_path = os.path.join(root, file)
            if any(os.path.abspath(os.path.join(root, d)) in ignore_dirs for d in dirs):
                continue
            content, original_path = handle_file(file_path)
            if content and not content.startswith("Unsupported"):
                file_index = write_to_output(content, output_dir, file_index, original_path)
            else:
                logging.info(content)

def limit_files_search(files, limit_search):
    if limit_search == 'noLimit':
        return files
    if limit_search == 'lastProducedPerType':
        file_types = {}
        for file in files:
            file_type = os.path.splitext(file)[1]
            if file_type not in file_types:
                file_types[file_type] = file
            else:
                if os.path.getmtime(file) > os.path.getmtime(file_types[file_type]):
                    file_types[file_type] = file
        return list(file_types.values())
    elif limit_search == 'lastProducedInFolder':
        if files:
            return [max(files, key=os.path.getmtime)]
    elif limit_search == 'lastProducedSimilarTitle':
        similar_titles = {}
        for file in files:
            base_name = os.path.splitext(file)[0]
            if base_name not in similar_titles:
                similar_titles[base_name] = file
            else:
                similarity = SequenceMatcher(None, base_name, os.path.splitext(similar_titles[base_name])[0]).ratio()
                if similarity > 0.9:
                    if os.path.getmtime(file) > os.path.getmtime(similar_titles[base_name]):
                        similar_titles[base_name] = file
        return list(similar_titles.values())
    return files

def handle_file(file_path):
    extension = os.path.splitext(file_path)[1].lower()
    handler = {
        '.txt': handle_text_file,
        '.htm': handle_epub_file,
        '.html': handle_epub_file,
        '.pdf': handle_pdf_file,
        '.docx': handle_word_file,
        '.doc': handle_word_file,
        '.pptx': handle_ppt_file,
        '.ppt': handle_ppt_file,
        '.xls': handle_excel_file,
        '.xlsx': handle_excel_file,
        '.xml': handle_xml_gan_file,
        '.gan': handle_xml_gan_file,
        '.xsd': handle_xml_gan_file,
        '.wav': handle_audio_file,
        '.mp3': handle_audio_file,
        '.m4a': handle_audio_file,
        '.mp4': handle_generic_video_file,
        '.avi': handle_generic_video_file,
        '.mov': handle_generic_video_file,
        '.mkv': handle_generic_video_file,
        '.mpeg': handle_generic_video_file,
        '.mpg': handle_generic_video_file,
        '.3gp': handle_generic_video_file,
        '.csv': handle_csv_file,
        '.epub': handle_epub_file,
        '.zip': handle_zip_file
    }.get(extension)
    if handler:
        return handler(file_path)
    return "Unsupported file format for {}".format(file_path), None

def process_text_with_keywords(text, keywords):
    json_data = []
    keyword_positions = []

    # Trova tutte le posizioni delle keyword nel testo
    for keyword in keywords:
        pattern = re.compile(keyword, re.IGNORECASE)
        matches = list(pattern.finditer(text))
        for match in matches:
            keyword_positions.append((match.start(), match.end(), match.group()))

    # Ordina le posizioni delle keyword in base alla loro posizione nel testo
    keyword_positions.sort()

    # Aggiungi i contenuti tra le keyword al JSON
    for i in range(len(keyword_positions)):
        start, end, matched_keyword = keyword_positions[i]
        next_start = keyword_positions[i + 1][0] if i + 1 < len(keyword_positions) else len(text)

        content = text[end:next_start].strip()
        json_data.append({"title": matched_keyword, "content": content})

    return json_data



def write_json(data, output_file):
    try:
        with open(output_file, 'w', encoding='utf-8') as json_file:
            json.dump(data, json_file, indent=4, ensure_ascii=False)
    except PermissionError:
        logging.error(f"Permission denied: {output_file}")
        messagebox.showerror("Error", f"Permission denied: {output_file}")
    except Exception as e:
        logging.error(f"Failed to write JSON file: {output_file} - {str(e)}")
        messagebox.showerror("Error", f"Failed to write JSON file: {output_file} - {str(e)}")

# Functions for GUI and video handling
def download_youtube_video(url, download_audio_only=False):
    if not url.strip():
        messagebox.showerror(lang['downloadError'], lang['downloadError'])
        return None
    try:
        yt = YouTube(url)
        title = yt.title
        title = ''.join([c for c in title if c.isalpha() or c.isdigit() or c==' ']).rstrip()
        if download_audio_only:
            stream = yt.streams.filter(only_audio=True).first()
            file_path = stream.download(output_path=temp_dir, filename=f"{title}.mp3")
        else:
            stream = yt.streams.get_highest_resolution()
            file_path = stream.download(output_path=temp_dir, filename=f"{title}.mp4")
        logging.info(f"{lang['downloadSuccess']}: {file_path}")
        return file_path
    except Exception as e:
        logging.error(f"{lang['downloadError']}: {e}")
        messagebox.showerror(lang['downloadError'], str(e))
        return None

def extract_audio(video_file):
    try:
        video = VideoFileClip(video_file)
        audio = video.audio
        audio_file = os.path.join(temp_dir, "temp_audio.wav")
        audio.write_audiofile(audio_file, codec='pcm_s16le')
        video.close()
        logging.info(f"Extracted audio to: {audio_file}")
        return audio_file
    except Exception as e:
        logging.error(f"{lang['audioExtractionError']}: {e}")
        messagebox.showerror(lang['audioExtractionError'], str(e))
        return None

def generate_srt(audio_file, output_file, language='it-IT'):
    recognizer = sr.Recognizer()
    sound = AudioSegment.from_wav(audio_file)
    chunks = split_on_silence(sound, min_silence_len=500, silence_thresh=sound.dBFS-14, keep_silence=500)

    with open(output_file, 'w') as file:
        start = 0
        for i, chunk in enumerate(chunks):
            chunk_filename = os.path.join(temp_dir, f"chunk{i}.wav")
            chunk.export(chunk_filename, format="wav")
            with sr.AudioFile(chunk_filename) as source:
                audio = recognizer.record(source)
            try:
                text = recognizer.recognize_google(audio, language=language)
                duration = len(chunk) / 1000
                start_time = start
                end_time = start + duration
                file.write(f"{i+1}\n")
                file.write(f"{format_time(start_time)} --> {format_time(end_time)}\n")
                file.write(f"{text.strip()}\n\n")
                start += duration
                logging.info(f"Generated SRT segment {i+1}")
            except sr.UnknownValueError:
                file.write(f"{i+1}\n")
                file.write(f"{format_time(start_time)} --> {format_time(end_time)}\n")
                file.write("Audio not understandable\n\n")
                start += duration
                logging.warning(f"Audio not understandable for segment {i+1}")
            except sr.RequestError as e:
                file.write(f"{i+1}\n")
                file.write(f"{format_time(start_time)} --> {format_time(end_time)}\n")
                file.write(f"Service error: {e}\n\n")
                start += duration
                logging.error(f"Service error for segment {i+1}: {e}")
    
    # Clean up chunk files
    for chunk_filename in os.listdir(temp_dir):
        if chunk_filename.startswith("chunk") and chunk_filename.endswith(".wav"):
            os.remove(os.path.join(temp_dir, chunk_filename))
            logging.info(f"Removed chunk file: {chunk_filename}")

def format_time(seconds):
    hours = int(seconds // 3600)
    minutes = int((seconds % 3600) // 60)
    seconds = seconds % 60
    return f"{hours:02}:{minutes:02}:{seconds:02},000"

def process_video():
    download_audio_only = download_audio_only_var.get()
    if url_entry.get().strip():
        video_file = download_youtube_video(url_entry.get().strip(), download_audio_only=download_audio_only)
    else:
        video_file = file_entry.get().strip()

    if video_file:
        if not download_audio_only:
            audio_file = extract_audio(video_file)
        else:
            audio_file = video_file

        if audio_file:
            output_path = filedialog.asksaveasfilename(defaultextension=".srt", filetypes=[("SRT files", "*.srt")])
            if output_path:
                generate_srt(audio_file, output_path, transcription_lang_var.get())
                messagebox.showinfo("Success", f"{lang['success']} {output_path}")
                logging.info(f"SRT file generated: {output_path}")

def setup_video_player(video_path, srt_path):
    media = instance.media_new(video_path)
    media.add_option(f'sub-file={srt_path}')
    player.set_media(media)
    player.play()

def load_video():
    video_path = video_entry.get().strip()
    srt_path = srt_entry.get().strip()

    if not video_path and not url_entry_video.get().strip():
        messagebox.showerror("Error", "Please enter a URL or select a video file.")
        return
    if not srt_path:
        messagebox.showerror("Error", "Please select an SRT file.")
        return

    if url_entry_video.get().strip():
        video_path = download_youtube_video(url_entry_video.get().strip())

    if video_path:
        setup_video_player(video_path, srt_path)

# Functions for Tab 3: Create text file
def start_create_text_file():
    directories = setup_directories
    output_path = setup_output_path
    ignore_dirs = setup_ignore_dirs
    process_subfolders = setup_process_subfolders
    limit_search = limit_search_var.get()
    for directory in directories:
        explore_directory(directory, output_path, ignore_dirs, process_subfolders, limit_search)
    logging.info(f"{lang['processCompleted']}: {directories}")

# Functions for Tab 4: Create Json
def start_create_json():
    directories = [setup_output_path]
    keywords = [entry.get() for entry in keyword_entries]
    combined_json_data = []

    for directory in directories:
        for root, dirs, files in os.walk(directory):
            for file in files:
                file_path = os.path.join(root, file)
                content, original_path = handle_file(file_path)
                if content and not content.startswith("Unsupported"):
                    json_data = process_text_with_keywords(content, keywords)
                    combined_json_data.append([
                        {"title": "Original file path", "content": original_path},
                        *json_data
                    ])
                else:
                    logging.info(content)

    output_path = setup_json_output_path
    if not os.path.exists(output_path):
        os.makedirs(output_path)
    json_output_file = os.path.join(output_path, "output.json")
    write_json(combined_json_data, json_output_file)
    messagebox.showinfo("Success", f"{lang['success']} {json_output_file}")
    logging.info(f"JSON file generated: {json_output_file}")
    
    # Aggiorna la configurazione dopo aver creato il JSON
    save_configuration()


# Functions for Tab 5: Setup
setup_directories = []
setup_ignore_dirs = []
setup_output_path = ""
setup_json_output_path = ""
setup_process_subfolders = True

keyword_entries = []

def add_keyword_entry():
    entry = tk.Entry(tab4, width=50)
    entry.pack(pady=5)
    keyword_entries.append(entry)

def remove_keyword_entry():
    if keyword_entries:
        entry = keyword_entries.pop()
        entry.pack_forget()

def add_directory():
    directory = filedialog.askdirectory()
    if directory:
        setup_directories.append(directory)
        directories_listbox.insert(tk.END, directory)

def remove_directory():
    selected = directories_listbox.curselection()
    if selected:
        index = selected[0]
        setup_directories.pop(index)
        directories_listbox.delete(index)

def add_ignore_directory():
    directory = filedialog.askdirectory()
    if directory:
        setup_ignore_dirs.append(directory)
        ignore_directories_listbox.insert(tk.END, directory)

def remove_ignore_directory():
    selected = ignore_directories_listbox.curselection()
    if selected:
        index = selected[0]
        setup_ignore_dirs.pop(index)
        ignore_directories_listbox.delete(index)

def set_output_path():
    global setup_output_path
    setup_output_path = filedialog.askdirectory()
    output_path_label.config(text=setup_output_path)

def set_json_output_path():
    global setup_json_output_path
    setup_json_output_path = filedialog.askdirectory()
    json_output_path_label.config(text=setup_json_output_path)

def set_process_subfolders():
    global setup_process_subfolders
    setup_process_subfolders = process_subfolders_var.get()

def set_temp_dir():
    global temp_dir
    temp_dir = filedialog.askdirectory()
    if not os.path.exists(temp_dir):
        os.makedirs(temp_dir)
    temp_dir_label.config(text=temp_dir)

# Save widget positions
def save_widget_positions():
    widget_positions = {}
    for widget in tab5.winfo_children():
        widget_positions[str(widget)] = {"x": widget.winfo_x(), "y": widget.winfo_y()}
    return widget_positions

# Load widget positions
def load_widget_positions(widget_positions):
    for widget in tab5.winfo_children():
        if str(widget) in widget_positions:
            pos = widget_positions[str(widget)]
            widget.place(x=pos["x"], y=pos["y"])

def update_setup_ui():
    directories_listbox.delete(0, tk.END)
    for directory in setup_directories:
        directories_listbox.insert(tk.END, directory)
    ignore_directories_listbox.delete(0, tk.END)
    for directory in setup_ignore_dirs:
        ignore_directories_listbox.insert(tk.END, directory)
    output_path_label.config(text=setup_output_path)
    json_output_path_label.config(text=setup_json_output_path)
    process_subfolders_var.set(setup_process_subfolders)
    temp_dir_label.config(text=temp_dir)
    limit_search_menu.set(lang[limit_search_var.get()])

def save_configuration():
    config = {
        "directories": setup_directories,
        "ignore_dirs": setup_ignore_dirs,
        "output_path": setup_output_path,
        "json_output_path": setup_json_output_path,
        "process_subfolders": setup_process_subfolders,
        "temp_dir": temp_dir,
        "limit_search": limit_search_var.get(),
        "keywords": [entry.get() for entry in keyword_entries],
        "widget_positions": save_widget_positions()
    }
    config_path = filedialog.asksaveasfilename(defaultextension=".json", filetypes=[("JSON files", "*.json")])
    if config_path:
        with open(config_path, 'w', encoding='utf-8') as config_file:
            json.dump(config, config_file, indent=4, ensure_ascii=False)
        messagebox.showinfo("Success", f"{lang['success']} {config_path}")
        logging.info(f"Configuration saved: {config_path}")

def load_configuration():
    config_path = filedialog.askopenfilename(filetypes=[("JSON files", "*.json")])
    if config_path:
        with open(config_path, 'r', encoding='utf-8') as config_file:
            config = json.load(config_file)
            global setup_directories, setup_ignore_dirs, setup_output_path, setup_json_output_path, setup_process_subfolders, temp_dir, limit_search_var
            setup_directories = config.get("directories", [])
            setup_ignore_dirs = config.get("ignore_dirs", [])
            setup_output_path = config.get("output_path", "")
            setup_json_output_path = config.get("json_output_path", "")
            setup_process_subfolders = config.get("process_subfolders", True)
            temp_dir = config.get("temp_dir", "temp")
            limit_search_var.set(config.get("limit_search", "noLimit"))
            
            keyword_entries.clear()
            for keyword in config.get("keywords", []):
                entry = tk.Entry(tab4, width=50)
                entry.insert(0, keyword)
                entry.pack(pady=5)
                keyword_entries.append(entry)
            
            if not os.path.exists(temp_dir):
                os.makedirs(temp_dir)
            update_setup_ui()
            load_widget_positions(config.get("widget_positions", {}))
        messagebox.showinfo("Success", f"{lang['success']} {config_path}")
        logging.info(f"Configuration loaded: {config_path}")

# GUI
root = TkinterDnD.Tk()
root.title("Magic a LoRA")

# Language selection
lang_var = tk.StringVar(root)
lang_var.set("en")  # default language
languages = {"en": "English", "it": "Italian", "fr": "French", "de": "German", "es": "Spanish", "ro": "Romanian", "pt": "Portuguese", "pl": "Polish", "sw": "Swahili"}
language_menu = tk.OptionMenu(root, lang_var, *languages.values(), command=lambda _: load_language(lang_var.get()))
language_menu.pack()

# Tabs
tab_control = ttk.Notebook(root)
tab1 = ttk.Frame(tab_control)
tab2 = ttk.Frame(tab_control)
tab3 = ttk.Frame(tab_control)
tab4 = ttk.Frame(tab_control)
tab5 = ttk.Frame(tab_control)
tab_control.add(tab1, text='Convert')
tab_control.add(tab2, text='Test SRT')
tab_control.add(tab3, text='Create text file')
tab_control.add(tab4, text='Create Json')
tab_control.add(tab5, text='Setup')
tab_control.pack(expand=1, fill='both')

# Tab 1: Create SRT
file_entry = tk.Entry(tab1, width=50)
file_entry.pack(pady=10)
browse_button = tk.Button(tab1, text=lang.get('browseVideo', 'Browse Video'), command=lambda: file_entry.insert(0, filedialog.askopenfilename()))
browse_button.pack(pady=10)

url_entry = tk.Entry(tab1, width=50)
url_entry.pack(pady=10)
url_label = tk.Label(tab1, text=lang.get('orEnterUrl', 'Or enter YouTube/Vimeo URL:'))
url_label.pack(pady=10)

download_audio_only_var = tk.BooleanVar()
download_audio_only_checkbox = tk.Checkbutton(tab1, text=lang.get('downloadAudioOnly', 'Download Audio Only'), variable=download_audio_only_var)
download_audio_only_checkbox.pack(pady=10)

transcription_lang_var = tk.StringVar(root)
transcription_lang_var.set("it-IT")
lang_options = ["en-US", "it-IT", "fr-FR", "de-DE", "es-ES", "pt-PT", "ro-RO", "pl-PL"]
transcription_lang_menu = tk.OptionMenu(tab1, transcription_lang_var, *lang_options)
transcription_lang_menu.pack(pady=10)

start_button = tk.Button(tab1, text=lang.get('startProcessing', 'Start Processing'), command=process_video)
start_button.pack(pady=20)

# Tab 2: Test SRT
video_entry = tk.Entry(tab2, width=50)
video_entry.pack(pady=10)
video_browse_button = tk.Button(tab2, text=lang.get('browseVideo', 'Browse Video'), command=lambda: video_entry.insert(0, filedialog.askopenfilename()))
video_browse_button.pack(pady=10)

url_entry_video = tk.Entry(tab2, width=50)
url_entry_video.pack(pady=10)
url_label_video = tk.Label(tab2, text=lang.get('orEnterUrlVideo', 'Or enter YouTube/Vimeo URL:'))
url_label_video.pack(pady=10)

srt_entry = tk.Entry(tab2, width=50)
srt_entry.pack(pady=10)
srt_browse_button = tk.Button(tab2, text=lang.get('browseSRTFile', 'Browse SRT File'), command=lambda: srt_entry.insert(0, filedialog.askopenfilename()))
srt_browse_button.pack(pady=10)

load_video_button = tk.Button(tab2, text=lang.get('loadAndPlayVideo', 'Load and Play Video'), command=lambda: threading.Thread(target=load_video).start())
load_video_button.pack(pady=20)

# VLC video frame
video_frame = ttk.Frame(tab2, height=400)
video_frame.pack(fill='both', expand=True)

# Initialize VLC player
instance = vlc.Instance()
player = instance.media_player_new()
player.set_hwnd(video_frame.winfo_id())

# Tab 3: Create text file
create_text_file_log = configure_logger('Create_text_file')
create_text_file_button = tk.Button(tab3, text=lang.get('startCreateTextFile', 'Start Create text file'), command=lambda: threading.Thread(target=start_create_text_file).start())
create_text_file_button.pack(pady=20)
create_text_file_log_display = tk.Text(tab3, height=15, state='disabled')
create_text_file_log_display.pack(fill='both', expand=True)

# Tab 4: Create Json
create_json_log = configure_logger('Create_Json')
chapter_keywords_label = tk.Label(tab4, text=lang.get('enterChapterKeywords', 'Enter chapter keywords (one per field):'))
chapter_keywords_label.pack(pady=10)
add_keyword_button = tk.Button(tab4, text=lang.get('addKeyword', 'Add Keyword'), command=add_keyword_entry)
add_keyword_button.pack(pady=5)
remove_keyword_button = tk.Button(tab4, text=lang.get('removeKeyword', 'Remove Keyword'), command=remove_keyword_entry)
remove_keyword_button.pack(pady=5)
create_json_button = tk.Button(tab4, text=lang.get('startCreateJson', 'Start Create Json'), command=lambda: threading.Thread(target=start_create_json).start())
create_json_button.pack(pady=20)
create_json_log_display = tk.Text(tab4, height=15, state='disabled')
create_json_log_display.pack(fill='both', expand=True)

# Tab 5: Setup
directories_listbox = tk.Listbox(tab5, selectmode=tk.SINGLE)
directories_listbox.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
add_directory_button = tk.Button(tab5, text=lang.get('addDirectory', 'Add Directory'), command=add_directory)
add_directory_button.pack(pady=5)
remove_directory_button = tk.Button(tab5, text=lang.get('removeDirectory', 'Remove Directory'), command=remove_directory)
remove_directory_button.pack(pady=5)

ignore_directories_listbox = tk.Listbox(tab5, selectmode=tk.SINGLE)
ignore_directories_listbox.pack(fill=tk.BOTH, expand=True, padx=10, pady=5)
add_ignore_directory_button = tk.Button(tab5, text=lang.get('addIgnoreDirectory', 'Add Ignore Directory'), command=add_ignore_directory)
add_ignore_directory_button.pack(pady=5)
remove_ignore_directory_button = tk.Button(tab5, text=lang.get('removeIgnoreDirectory', 'Remove Ignore Directory'), command=remove_ignore_directory)
remove_ignore_directory_button.pack(pady=5)

output_path_label = tk.Label(tab5, text="")
output_path_label.pack(pady=10)
set_output_path_button = tk.Button(tab5, text=lang.get('setTextOutput', 'Set Text Output'), command=set_output_path)
set_output_path_button.pack(pady=5)

json_output_path_label = tk.Label(tab5, text="")
json_output_path_label.pack(pady=10)
set_json_output_path_button = tk.Button(tab5, text=lang.get('setJsonOutput', 'Set JSON Output'), command=set_json_output_path)
set_json_output_path_button.pack(pady=5)

temp_dir_label = tk.Label(tab5, text=temp_dir)
temp_dir_label.pack(pady=10)
set_temp_dir_button = tk.Button(tab5, text=lang.get('setTempDir', 'Set Temp Directory'), command=set_temp_dir)
set_temp_dir_button.pack(pady=5)

limit_search_var = tk.StringVar()
limit_search_label = tk.Label(tab5, text=lang.get('limitSearch', 'Limit Search'))
limit_search_label.pack(pady=10)
limit_search_menu = ttk.Combobox(tab5, textvariable=limit_search_var)
limit_search_menu['values'] = [lang.get('noLimit', 'No Limit'), lang.get('lastProducedPerType', 'Last Produced per Type'), lang.get('lastProducedInFolder', 'Last Produced in Folder'), lang.get('lastProducedSimilarTitle', 'Last Produced with Similar Title')]
limit_search_menu.set(lang.get('noLimit', 'No Limit'))
limit_search_menu.pack(pady=10)

process_subfolders_var = tk.BooleanVar()
process_subfolders_checkbox = tk.Checkbutton(tab5, text=lang.get('processSubfolders', 'Process Subfolders'), variable=process_subfolders_var, command=set_process_subfolders)
process_subfolders_checkbox.pack(pady=5)

save_config_button = tk.Button(tab5, text=lang.get('saveConfiguration', 'Save Configuration'), command=save_configuration)
save_config_button.pack(pady=5)
load_config_button = tk.Button(tab5, text=lang.get('loadConfiguration', 'Load Configuration'), command=load_configuration)
load_config_button.pack(pady=5)

# Enable dragging for the widgets in the Setup tab
def make_draggable(widget):
    widget.bind("<Button-1>", on_drag_start)
    widget.bind("<B1-Motion>", on_drag_motion)

def on_drag_start(event):
    widget = event.widget
    widget._drag_start_x = event.x
    widget._drag_start_y = event.y

def on_drag_motion(event):
    widget = event.widget
    x = widget.winfo_x() - widget._drag_start_x + event.x
    y = widget.winfo_y() - widget._drag_start_y + event.y
    widget.place(x=x, y=y)

make_draggable(directories_listbox)
make_draggable(add_directory_button)
make_draggable(remove_directory_button)
make_draggable(ignore_directories_listbox)
make_draggable(add_ignore_directory_button)
make_draggable(remove_ignore_directory_button)
make_draggable(output_path_label)
make_draggable(set_output_path_button)
make_draggable(json_output_path_label)
make_draggable(set_json_output_path_button)
make_draggable(temp_dir_label)
make_draggable(set_temp_dir_button)
make_draggable(limit_search_label)
make_draggable(limit_search_menu)
make_draggable(process_subfolders_checkbox)
make_draggable(save_config_button)
make_draggable(load_config_button)

def on_closing():
    for thread in threading.enumerate():
        if thread is not threading.main_thread():
            thread.join(timeout=1)
    root.destroy()

root.protocol("WM_DELETE_WINDOW", on_closing)
root.mainloop()
